#!/usr/bin/env python3
"""
generate_pptx.py
Minimal startup-style PowerPoint dla projektu GGSN Temporal GNN.

Uruchomienie:
    uv run python generate_pptx.py

Wyjście: presentation.pptx
"""

from __future__ import annotations

import io
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.patches as FancyPatch
import matplotlib.patches as mpatches
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).parent / "presentation.pptx"

# ─── Color palette (minimal white / startup) ──────────────────────────────────
W, H = Inches(13.33), Inches(7.5)  # 16:9


def rgb(hex_: str) -> RGBColor:
    h = hex_.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


C_BG = rgb("FFFFFF")
C_INK = rgb("0B0F1E")  # NeuroNauts: very dark navy titles
C_TEXT = rgb("1F293A")  # NeuroNauts: dark body
C_MUTED = rgb("64748B")  # secondary text
C_LIGHT = rgb("94A3B8")  # captions
C_RULE = rgb("E2E8F0")  # horizontal rules
C_CARD = rgb("EEECFB")  # NeuroNauts: lavender card bg (signature)
C_BORDER = rgb("E2E8F0")  # card border
C_ACCENT = rgb("4F46E5")  # NeuroNauts: indigo-600
C_GREEN = rgb("10B981")  # emerald (positive outcomes)
C_RED = rgb("EF4444")  # red (negative / ablation)
C_AMBER = rgb("F59E0B")  # amber (best model highlight)
C_CYAN = rgb("06B6D4")  # cyan (secondary)

FONT = "Inter"

# ─── Experiment data (from data/snapshots/gnn/) ───────────────────────────────
EXPS = [
    ("signal only", 0.784, 0.381, "ablation"),
    ("e2e fine-tuning", 0.790, 0.303, "ablation"),
    ("+ allstay (first-50)", 0.820, 0.396, "feature"),
    ("+ ICD node", 0.823, 0.377, "feature"),
    ("+ focal γ=2 only", 0.827, 0.408, "loss"),
    ("baseline", 0.829, 0.403, "baseline"),
    ("2 GIN layers", 0.835, 0.438, "arch"),
    ("demo + 4 layers", 0.836, 0.408, "arch"),
    ("+ demographics", 0.841, 0.429, "demo"),
    ("demo + focal γ=2", 0.843, 0.426, "combined"),
    ("demo + attention pool", 0.844, 0.448, "pooling"),
    ("demo + dual pool", 0.846, 0.435, "pooling"),
    ("demo + attn + focal (best)", 0.850, 0.465, "best"),
]

GROUP_COLORS_LIGHT = {
    "ablation": "#CBD5E1",  # slate-300  — nie pomaga
    "feature": "#A5B4FC",  # indigo-300
    "baseline": "#818CF8",  # indigo-400
    "loss": "#6366F1",  # indigo-500
    "arch": "#4F46E5",  # indigo-600 (accent)
    "demo": "#4338CA",  # indigo-700
    "combined": "#3730A3",  # indigo-800
    "pooling": "#312E81",  # indigo-900
    "best": "#F59E0B",  # amber — wyróżniony złotem
}


# ─── Matplotlib light theme ────────────────────────────────────────────────────
plt.rcParams.update(
    {
        "figure.facecolor": "white",
        "axes.facecolor": "white",
        "axes.edgecolor": "#E2E8F0",
        "axes.labelcolor": "#475569",
        "xtick.color": "#94A3B8",
        "ytick.color": "#475569",
        "text.color": "#0F172A",
        "grid.color": "#F1F5F9",
        "grid.linewidth": 0.8,
        "font.family": "DejaVu Sans",
        "font.size": 11,
        "axes.spines.top": False,
        "axes.spines.right": False,
    }
)


def _png(fig) -> io.BytesIO:
    buf = io.BytesIO()
    fig.savefig(
        buf, format="png", bbox_inches="tight", facecolor="white", edgecolor="none", dpi=180
    )
    buf.seek(0)
    plt.close(fig)
    return buf


# ─── Charts ───────────────────────────────────────────────────────────────────


def chart_auroc() -> io.BytesIO:
    names = [e[0] for e in EXPS]
    aurocs = [e[1] for e in EXPS]
    groups = [e[3] for e in EXPS]
    colors = [GROUP_COLORS_LIGHT[g] for g in groups]

    fig, ax = plt.subplots(figsize=(10, 4.6))
    y = np.arange(len(names))
    bars = ax.barh(y, [a - 0.74 for a in aurocs], left=0.74, color=colors, height=0.62, zorder=3)

    ax.axvline(0.88, color="#64748B", lw=1.5, ls="--", alpha=0.7, zorder=4)
    ax.text(0.8815, len(names) - 0.6, "SOTA 0.88", color="#64748B", fontsize=8.5, va="top")
    ax.axvline(0.829, color="#6366F1", lw=1, ls=":", alpha=0.45, zorder=4)

    for bar, v, g in zip(bars, aurocs, groups):
        fw = "bold" if g in ("best", "baseline") else "normal"
        fc = "#1E40AF" if g == "best" else "#0F172A"
        ax.text(
            v + 0.0012,
            bar.get_y() + bar.get_height() / 2,
            f"{v:.3f}",
            va="center",
            ha="left",
            fontsize=9.5,
            fontweight=fw,
            color=fc,
        )

    ax.set_yticks(y)
    ax.set_yticklabels(names, fontsize=10)
    ax.set_xlim(0.74, 0.895)
    ax.set_xlabel("AUROC", fontsize=11, color="#475569")
    ax.set_title(
        "AUROC — 13 eksperymentów",
        fontsize=13,
        pad=12,
        fontweight="bold",
        color="#0F172A",
        loc="left",
    )
    ax.grid(axis="x", zorder=0)
    ax.set_axisbelow(True)

    seen = {}
    for g, c in GROUP_COLORS_LIGHT.items():
        if g in groups and g not in seen:
            seen[g] = c
    patches = [mpatches.Patch(color=c, label=g) for g, c in seen.items()]
    ax.legend(
        handles=patches,
        fontsize=8,
        loc="lower right",
        facecolor="white",
        edgecolor="#E2E8F0",
        labelcolor="#475569",
        framealpha=1,
    )
    fig.tight_layout()
    return _png(fig)


def chart_auprc() -> io.BytesIO:
    names = [e[0] for e in EXPS]
    auprcs = [e[2] for e in EXPS]
    groups = [e[3] for e in EXPS]
    colors = [GROUP_COLORS_LIGHT[g] for g in groups]

    fig, ax = plt.subplots(figsize=(10, 4.6))
    y = np.arange(len(names))
    ax.barh(y, [a - 0.25 for a in auprcs], left=0.25, color=colors, height=0.62, zorder=3)

    ax.axvline(0.127, color="#1E40AF", lw=1.5, ls="--", alpha=0.7, zorder=4)
    ax.text(0.134, len(names) - 0.6, "random 0.127", color="#1E40AF", fontsize=8.5, va="top")

    for i, (v, g) in enumerate(zip(auprcs, groups)):
        fc = "#1E40AF" if g == "best" else "#0F172A"
        fw = "bold" if g == "best" else "normal"
        ax.text(
            v + 0.004,
            i + 0,
            f"{v:.3f}",
            va="center",
            ha="left",
            fontsize=9.5,
            fontweight=fw,
            color=fc,
        )

    ax.set_yticks(y)
    ax.set_yticklabels(names, fontsize=10)
    ax.set_xlim(0.25, 0.53)
    ax.set_xlabel("AUPRC", fontsize=11, color="#475569")
    ax.set_title(
        "AUPRC — precision–recall  (random baseline 0.127)",
        fontsize=13,
        pad=12,
        fontweight="bold",
        color="#0F172A",
        loc="left",
    )
    ax.grid(axis="x", zorder=0)
    ax.set_axisbelow(True)
    fig.tight_layout()
    return _png(fig)


def chart_leakage() -> io.BytesIO:
    fig, ax = plt.subplots(figsize=(6, 3.6))
    labels = ["ostatnie 50 sygnałów\n(temporal leakage)", "pierwsze 50 sygnałów\n(naprawione)"]
    vals = [0.931, 0.820]
    colors = ["#64748B", "#2563EB"]

    bars = ax.bar([0, 1], vals, color=colors, width=0.5, zorder=3)
    ax.set_ylim(0.75, 0.97)
    ax.set_ylabel("AUROC", color="#475569")
    ax.set_xticks([0, 1])
    ax.set_xticklabels(labels, fontsize=10)
    ax.set_title(
        "Temporal leakage bug", fontsize=13, pad=12, fontweight="bold", color="#0F172A", loc="left"
    )

    for bar, v, c in zip(bars, vals, colors):
        ax.text(
            bar.get_x() + bar.get_width() / 2,
            v + 0.005,
            f"{v:.3f}",
            ha="center",
            fontsize=16,
            fontweight="bold",
            color=c,
        )

    ax.annotate(
        "sygnały tuż przed\nśmiercią = leakage",
        xy=(0, 0.931),
        xytext=(0.52, 0.890),
        arrowprops=dict(arrowstyle="->", color="#64748B", lw=1.4),
        color="#64748B",
        fontsize=9,
    )
    ax.grid(axis="y", zorder=0)
    ax.set_axisbelow(True)
    fig.tight_layout()
    return _png(fig)


def chart_contributions() -> io.BytesIO:
    baseline = 0.829
    components = [
        ("− notatki kliniczne", 0.784 - baseline, "#64748B"),
        ("+ ICD node", 0.823 - baseline, "#64748B"),
        ("+ allstay first-50", 0.820 - baseline, "#64748B"),
        ("+ focal γ=2 only", 0.827 - baseline, "#3B82F6"),
        ("+ demographics", 0.841 - baseline, "#2563EB"),
        ("+ attention pooling", 0.844 - baseline, "#06B6D4"),
        ("+ dual pooling", 0.846 - baseline, "#2563EB"),
        ("demo + attn + focal (best)", 0.850 - baseline, "#1E40AF"),
    ]
    labels = [c[0] for c in components]
    deltas = [c[1] for c in components]
    colors = [c[2] for c in components]

    fig, ax = plt.subplots(figsize=(7.5, 4.2))
    ax.barh(range(len(labels)), deltas, color=colors, height=0.58, zorder=3)
    ax.axvline(0, color="#94A3B8", lw=1.5, zorder=4)

    for i, (d, c) in enumerate(zip(deltas, colors)):
        off = 0.0008 if d >= 0 else -0.0008
        ha = "left" if d >= 0 else "right"
        ax.text(
            d + off, i, f"{d:+.3f}", va="center", ha=ha, fontsize=10, color=c, fontweight="bold"
        )

    ax.set_yticks(range(len(labels)))
    ax.set_yticklabels(labels, fontsize=10.5)
    ax.set_xlabel("ΔAUROC vs baseline (0.829)", fontsize=11, color="#475569")
    ax.set_title(
        "Wkład komponentów", fontsize=13, pad=12, fontweight="bold", color="#0F172A", loc="left"
    )
    ax.grid(axis="x", zorder=0)
    ax.set_axisbelow(True)
    fig.tight_layout()
    return _png(fig)


def chart_sota() -> io.BytesIO:
    fig, ax = plt.subplots(figsize=(6, 3.6))
    models = ["Random", "Signal\nonly", "Nasz\nbest", "SOTA"]
    vals = [0.50, 0.784, 0.850, 0.880]
    colors = ["#CBD5E1", "#64748B", "#1E40AF", "#6366F1"]

    bars = ax.bar(range(4), vals, color=colors, width=0.55, zorder=3)
    ax.set_ylim(0.40, 0.930)
    ax.set_ylabel("AUROC", color="#475569")
    ax.set_xticks(range(4))
    ax.set_xticklabels(models, fontsize=11)
    ax.set_title(
        "Nasz wynik vs SOTA", fontsize=13, pad=12, fontweight="bold", color="#0F172A", loc="left"
    )

    for bar, v, c in zip(bars, vals, colors):
        ax.text(
            bar.get_x() + bar.get_width() / 2,
            v + 0.012,
            f"{v:.3f}",
            ha="center",
            fontsize=13,
            fontweight="bold",
            color=c,
        )

    ax.annotate(
        "",
        xy=(3, 0.880),
        xytext=(2, 0.850),
        arrowprops=dict(arrowstyle="<->", color="#94A3B8", lw=1.5),
    )
    ax.text(2.57, 0.875, "Δ 0.030", ha="center", color="#475569", fontsize=10)
    ax.grid(axis="y", zorder=0)
    ax.set_axisbelow(True)
    fig.tight_layout()
    return _png(fig)


def chart_learning_curve() -> io.BytesIO:
    import csv

    path = Path(__file__).parent / "data" / "snapshots" / "gnn" / "version_14" / "metrics.csv"
    rows = list(csv.DictReader(path.open()))

    val_epochs, val_auroc = [], []
    train_epochs, train_loss = [], []
    for r in rows:
        ep = int(r["epoch"])
        if r.get("val_auroc"):
            val_epochs.append(ep)
            val_auroc.append(float(r["val_auroc"]))
        if r.get("train_loss"):
            train_epochs.append(ep)
            train_loss.append(float(r["train_loss"]))

    fig, ax1 = plt.subplots(figsize=(10, 4.6))
    ax2 = ax1.twinx()

    ax1.plot(
        val_epochs,
        val_auroc,
        color="#6366F1",
        lw=2.2,
        zorder=3,
        marker="o",
        markersize=3.5,
        label="val AUROC",
    )
    best_idx = int(np.argmax(val_auroc))
    ax1.scatter(
        [val_epochs[best_idx]],
        [val_auroc[best_idx]],
        s=120,
        color="#1E40AF",
        zorder=5,
        edgecolor="white",
        lw=1.5,
    )
    ax1.annotate(
        f"best 0.{int(val_auroc[best_idx] * 1000):03d} @ ep.{val_epochs[best_idx]}",
        xy=(val_epochs[best_idx], val_auroc[best_idx]),
        xytext=(val_epochs[best_idx] - 6, val_auroc[best_idx] - 0.025),
        fontsize=10.5,
        color="#1E3A8A",
        fontweight="bold",
    )

    ax2.plot(
        train_epochs, train_loss, color="#94A3B8", lw=1.3, alpha=0.6, zorder=2, label="train loss"
    )

    ax1.set_xlabel("epoka", fontsize=11, color="#475569")
    ax1.set_ylabel("val AUROC", fontsize=11, color="#6366F1")
    ax2.set_ylabel("train loss (focal)", fontsize=10.5, color="#94A3B8")
    ax1.set_ylim(0.75, 0.86)
    ax1.set_title(
        "Krzywa uczenia — demo_attention_focal2 (version_14)",
        fontsize=13,
        pad=12,
        fontweight="bold",
        color="#0F172A",
        loc="left",
    )
    ax1.grid(axis="y", zorder=0, alpha=0.5)
    ax1.set_axisbelow(True)
    ax1.tick_params(axis="y", labelcolor="#6366F1")
    ax2.tick_params(axis="y", labelcolor="#94A3B8")

    h1, l1 = ax1.get_legend_handles_labels()
    h2, l2 = ax2.get_legend_handles_labels()
    ax1.legend(
        h1 + h2,
        l1 + l2,
        fontsize=9.5,
        loc="lower right",
        facecolor="white",
        edgecolor="#E2E8F0",
        framealpha=1,
    )

    fig.tight_layout()
    return _png(fig)


def chart_similarity() -> io.BytesIO:
    """
    50×50 cosine similarity matrix: 10 stays × 5 notatek per stay.
    Block-diagonal pattern = contrastive nauczył kontekstu pobytu.
    """
    import pandas as pd

    root = Path(__file__).parent
    cache = root / "data" / "processed" / "simmat_cache.npz"
    emb_path = root / "data" / "embeddings" / "node_embeddings.pt"
    pairs_path = root / "data" / "processed" / "pairs_all-icus_note_level.csv"

    N_STAYS = 10
    N_NOTES_PER_STAY = 5

    if cache.exists():
        d = np.load(cache)
        sim = d["sim"]
        same_stay_sims = d["same"]
        diff_stay_sims = d["diff"]
    else:
        print("  [simmat] budowanie cache (jednorazowe, ~1-2 min)...")
        emb = torch.load(emb_path, weights_only=False)

        note_to_stay: dict[str, str] = {}
        for chunk in pd.read_csv(pairs_path, usecols=["note_id", "stay_id"], chunksize=2_000_000):
            for nid, sid in zip(chunk["note_id"].values, chunk["stay_id"].values):
                if nid not in note_to_stay:
                    note_to_stay[nid] = str(sid)

        from collections import defaultdict

        stay_to_notes: dict[str, list[str]] = defaultdict(list)
        for nid in emb.keys():
            sid = note_to_stay.get(nid)
            if sid is not None:
                stay_to_notes[sid].append(nid)
        eligible = [sid for sid, ns in stay_to_notes.items() if len(ns) >= N_NOTES_PER_STAY]
        rng = np.random.default_rng(42)
        picked_stays = rng.choice(eligible, N_STAYS, replace=False)

        selected_ids: list[str] = []
        for sid in picked_stays:
            notes = sorted(stay_to_notes[sid])[:N_NOTES_PER_STAY]
            selected_ids.extend(notes)

        X = np.stack([emb[nid].cpu().numpy() for nid in selected_ids])
        X = X / (np.linalg.norm(X, axis=1, keepdims=True) + 1e-9)
        sim = X @ X.T  # 50×50 cosine (embeddingi już L2-normalized w Phase 1)

        same_mask = np.zeros_like(sim, dtype=bool)
        for k in range(N_STAYS):
            s = k * N_NOTES_PER_STAY
            same_mask[s : s + N_NOTES_PER_STAY, s : s + N_NOTES_PER_STAY] = True
        np.fill_diagonal(same_mask, False)
        diff_mask = ~same_mask
        np.fill_diagonal(diff_mask, False)

        same_stay_sims = sim[same_mask]
        diff_stay_sims = sim[diff_mask]
        np.savez(cache, sim=sim, same=same_stay_sims, diff=diff_stay_sims)

    fig, (ax1, ax2) = plt.subplots(
        1, 2, figsize=(11, 4.6), gridspec_kw={"width_ratios": [1.0, 1.1]}
    )

    from matplotlib.colors import LinearSegmentedColormap

    blues_cmap = LinearSegmentedColormap.from_list(
        "blues_custom", ["#F8FAFC", "#DBEAFE", "#93C5FD", "#3B82F6", "#1E3A8A"]
    )

    im = ax1.imshow(
        sim, cmap=blues_cmap, vmin=-0.2, vmax=1.0, aspect="equal", interpolation="nearest"
    )
    for k in range(1, N_STAYS):
        ax1.axhline(k * N_NOTES_PER_STAY - 0.5, color="white", lw=0.8)
        ax1.axvline(k * N_NOTES_PER_STAY - 0.5, color="white", lw=0.8)

    ax1.set_xticks([])
    ax1.set_yticks([])
    ax1.set_xlabel(
        f"{N_STAYS * N_NOTES_PER_STAY} notatek ({N_STAYS} pobytów × {N_NOTES_PER_STAY})",
        fontsize=10,
        color="#475569",
    )
    ax1.set_title(
        "Macierz podobieństwa  (cosine)",
        fontsize=12,
        pad=10,
        fontweight="bold",
        color="#0F172A",
        loc="left",
    )
    cb = fig.colorbar(im, ax=ax1, fraction=0.046, pad=0.04)
    cb.outline.set_visible(False)
    cb.ax.tick_params(labelsize=8, colors="#64748B")

    bins = np.linspace(-0.2, 1.0, 40)
    ax2.hist(
        diff_stay_sims,
        bins=bins,
        color="#CBD5E1",
        alpha=0.85,
        label=f"różne pobyty  (n={len(diff_stay_sims)})",
        zorder=2,
    )
    ax2.hist(
        same_stay_sims,
        bins=bins,
        color="#1E3A8A",
        alpha=0.85,
        label=f"ten sam pobyt  (n={len(same_stay_sims)})",
        zorder=3,
    )
    ax2.axvline(diff_stay_sims.mean(), color="#64748B", ls="--", lw=1.2)
    ax2.axvline(same_stay_sims.mean(), color="#1E3A8A", ls="--", lw=1.2)
    ax2.text(
        diff_stay_sims.mean(),
        ax2.get_ylim()[1] * 0.95,
        f" śr {diff_stay_sims.mean():.2f}",
        color="#64748B",
        fontsize=9,
        va="top",
    )
    ax2.text(
        same_stay_sims.mean(),
        ax2.get_ylim()[1] * 0.85,
        f" śr {same_stay_sims.mean():.2f}",
        color="#1E3A8A",
        fontsize=9,
        va="top",
        fontweight="bold",
    )
    ax2.set_xlabel("cosine similarity", fontsize=10, color="#475569")
    ax2.set_ylabel("liczność", fontsize=10, color="#475569")
    ax2.set_title(
        "Rozkład podobieństw", fontsize=12, pad=10, fontweight="bold", color="#0F172A", loc="left"
    )
    ax2.legend(
        fontsize=9, loc="upper right", facecolor="white", edgecolor="#E2E8F0", framealpha=0.95
    )

    fig.tight_layout()
    return _png(fig)


import torch  # noqa: E402  (used by chart_similarity / chart_umap)


def chart_umap() -> io.BytesIO:
    """PCA-2D projekcja embeddingów notatek po Phase 1, kolorowana mortality."""
    from sklearn.decomposition import PCA

    root = Path(__file__).parent
    emb_path = root / "data" / "embeddings" / "node_embeddings.pt"
    pairs_path = root / "data" / "processed" / "pairs_all-icus_note_level.csv"
    cache = root / "data" / "processed" / "pca2d_cache.npz"

    if cache.exists():
        d = np.load(cache)
        X2, y = d["X2"], d["y"]
    else:
        print("  [chart_umap] ladowanie embeddingów (jednorazowe ~15s)…")
        import csv

        note_mort: dict[str, int] = {}
        with open(pairs_path) as f:
            for row in csv.DictReader(f):
                nid = row["note_id"]
                if nid not in note_mort:
                    note_mort[nid] = int(float(row["mortality"]))

        emb_dict = torch.load(emb_path, map_location="cpu", weights_only=False)
        available = [nid for nid in emb_dict if nid in note_mort]

        rng = np.random.default_rng(42)
        n = min(3000, len(available))
        sampled = rng.choice(available, n, replace=False)

        X = np.stack([emb_dict[nid].cpu().numpy() for nid in sampled])
        y = np.array([note_mort[nid] for nid in sampled], dtype=np.int8)

        X2 = PCA(n_components=2, random_state=42).fit_transform(X)
        np.savez(cache, X2=X2, y=y)
        print(f"     cache zapisany: {cache}")

    fig, ax = plt.subplots(figsize=(8.5, 5.0))
    mask0, mask1 = y == 0, y == 1
    ax.scatter(
        X2[mask0, 0],
        X2[mask0, 1],
        c="#93C5FD",
        s=4,
        alpha=0.22,
        label=f"przezycie  (n={mask0.sum():,})",
        rasterized=True,
    )
    ax.scatter(
        X2[mask1, 0],
        X2[mask1, 1],
        c="#F87171",
        s=12,
        alpha=0.65,
        label=f"zgon  (n={mask1.sum():,})",
        rasterized=True,
    )

    ax.set_title(
        "PCA-2D embeddingów notatek (Phase 1)  —  bez fine-tuning na mortality",
        fontsize=12,
        pad=10,
        fontweight="bold",
        color="#0F172A",
        loc="left",
    )
    ax.set_xlabel("PC1", fontsize=10, color="#475569")
    ax.set_ylabel("PC2", fontsize=10, color="#475569")
    ax.legend(fontsize=10, facecolor="white", edgecolor="#E2E8F0", markerscale=3, framealpha=1)
    ax.grid(alpha=0.3)
    ax.set_axisbelow(True)
    fig.tight_layout()
    return _png(fig)


# ─── Diagram figures ──────────────────────────────────────────────────────────


def diagram_twotower() -> io.BytesIO:
    """
    Layout (y top→bottom, data coords 0–6):
      5.6-6.0  title
      4.5-5.4  Input boxes
      4.5-3.9  arrows down
      2.2-3.9  Tower boxes
      2.2-1.8  arrows to embeddings
      1.6      embedding labels
      1.6-1.1  arrows →  InfoNCE  ←
      0.2-1.1  InfoNCE box   ← clearly below towers, 1.1 gap
    """
    fig, ax = plt.subplots(figsize=(11, 5.2))
    ax.set_xlim(0, 11)
    ax.set_ylim(0, 6.0)
    ax.axis("off")

    ACC = "#4F46E5"
    ACC_L = "#EEECFB"
    ACC_B = "#C7D2FE"

    def box(x, y, w, h, fc, ec, title, sub, tc=None):
        p = FancyPatch.FancyBboxPatch(
            (x, y),
            w,
            h,
            boxstyle="round,pad=0.07",
            facecolor=fc,
            edgecolor=ec,
            linewidth=1.5,
            zorder=3,
        )
        ax.add_patch(p)
        tc = tc or "#0B0F1E"
        ax.text(
            x + w / 2,
            y + h - 0.28,
            title,
            ha="center",
            va="top",
            fontsize=11,
            fontweight="bold",
            color=tc,
            zorder=4,
        )
        if sub:
            ax.text(
                x + w / 2,
                y + 0.55,
                sub,
                ha="center",
                va="center",
                fontsize=8.5,
                color="#64748B",
                zorder=4,
                linespacing=1.5,
            )

    def arr(x1, y1, x2, y2, color="#94A3B8", lw=1.5):
        ax.annotate(
            "",
            xy=(x2, y2),
            xytext=(x1, y1),
            arrowprops=dict(arrowstyle="-|>", color=color, lw=lw, mutation_scale=14),
        )

    # ── Inputs (top)
    box(
        0.3,
        4.5,
        3.2,
        0.9,
        "#F0F9FF",
        "#BAE6FD",
        "Notatka kliniczna",
        '"Patient appears acutely ill..."',
        "#0369A1",
    )
    box(7.5, 4.5, 3.2, 0.9, ACC_L, ACC_B, "Sygnaly fizjologiczne", "HR=82, SpO2=97, MAP=68", ACC)

    # ── Arrows input → tower
    arr(1.9, 4.5, 1.9, 3.92)
    arr(9.1, 4.5, 9.1, 3.92)

    # ── Tower boxes (generous height)
    box(
        0.3,
        2.1,
        3.2,
        1.75,
        ACC_L,
        ACC_B,
        "TextTower",
        "Bio_ClinicalBERT\n768 → 384 → 128\nL2-normalize",
        ACC,
    )
    box(
        7.5,
        2.1,
        3.2,
        1.75,
        ACC_L,
        ACC_B,
        "SignalTower",
        "Embedding(14,8) + 1D-CNN\nGAP → Linear(128)\nL2-normalize",
        ACC,
    )

    # ── Arrows tower → embedding labels
    arr(1.9, 2.1, 1.9, 1.75)
    arr(9.1, 2.1, 9.1, 1.75)

    # ── Embedding labels
    ax.text(1.9, 1.55, "z_text  ∈  ℝ¹²⁸", ha="center", fontsize=10.5, fontweight="bold", color=ACC)
    ax.text(
        9.1, 1.55, "z_signal  ∈  ℝ¹²⁸", ha="center", fontsize=10.5, fontweight="bold", color=ACC
    )

    # ── Arrows → InfoNCE (clearly in the space BELOW towers)
    arr(3.5, 1.55, 4.3, 1.25, color=ACC, lw=2.0)
    arr(7.5, 1.55, 6.7, 1.25, color=ACC, lw=2.0)

    # ── InfoNCE box — centred, well below towers (h=1.1 so title+sub fit)
    box(4.2, 0.15, 2.6, 1.1, ACC_L, ACC, "InfoNCE  (NT-Xent)", "τ = 0.07  ·  hard negatives", ACC)

    # ── Shared embedding space label (between towers)
    ax.text(
        5.5,
        2.9,
        "wspólna przestrzeń 128-dim",
        ha="center",
        fontsize=10,
        color=ACC,
        style="italic",
        zorder=5,
        bbox=dict(boxstyle="round,pad=0.2", fc="white", ec=ACC_B, lw=0.8),
    )

    # ── Title
    ax.text(
        0.3,
        5.85,
        "PHASE 1 — Contrastive Pre-training",
        fontsize=12,
        fontweight="bold",
        color="#0B0F1E",
    )

    fig.tight_layout(pad=0.3)
    return _png(fig)


def diagram_graph() -> io.BytesIO:
    fig, ax = plt.subplots(figsize=(11, 3.8))
    ax.set_xlim(0, 11)
    ax.set_ylim(-0.3, 3.8)
    ax.axis("off")

    # Timeline
    ax.plot([0.3, 10.7], [0.5, 0.5], color="#E2E8F0", lw=1.5, zorder=0)
    for tx, label in [(0.8, "t = −1h"), (2.5, "0h"), (4.5, "6h"), (6.5, "14h"), (8.8, "22h")]:
        ax.text(tx, 0.15, label, ha="center", fontsize=9, color="#94A3B8")
        ax.plot([tx, tx], [0.4, 0.6], color="#CBD5E1", lw=1)

    def node_sig(x, label, val):
        circ = plt.Circle((x, 1.8), 0.52, color="#06B6D4", alpha=0.15, zorder=3)
        ax.add_patch(circ)
        circ2 = plt.Circle((x, 1.8), 0.52, fill=False, ec="#06B6D4", lw=1.5, zorder=4)
        ax.add_patch(circ2)
        ax.text(
            x,
            1.95,
            label,
            ha="center",
            va="center",
            fontsize=9,
            fontweight="bold",
            color="#0369A1",
            zorder=5,
        )
        ax.text(x, 1.62, val, ha="center", va="center", fontsize=8, color="#64748B", zorder=5)

    def node_note(x, label):
        rect = FancyPatch.FancyBboxPatch(
            (x - 0.6, 1.2),
            1.2,
            0.7,
            boxstyle="round,pad=0.05",
            facecolor="#EEF2FF",
            edgecolor="#818CF8",
            linewidth=1.5,
            zorder=3,
        )
        ax.add_patch(rect)
        ax.text(
            x,
            1.62,
            label,
            ha="center",
            va="center",
            fontsize=9,
            fontweight="bold",
            color="#4338CA",
            zorder=4,
        )
        ax.text(x, 1.38, "128-dim", ha="center", va="center", fontsize=8, color="#64748B", zorder=4)

    def node_icd(x):
        diamond = plt.Polygon(
            [(x, 2.35), (x + 0.5, 1.8), (x, 1.25), (x - 0.5, 1.8)],
            facecolor="#EEF2FF",
            edgecolor="#818CF8",
            linewidth=1.5,
            zorder=3,
        )
        ax.add_patch(diamond)
        ax.text(
            x,
            1.88,
            "ICD",
            ha="center",
            va="center",
            fontsize=9,
            fontweight="bold",
            color="#4338CA",
            zorder=4,
        )
        ax.text(
            x, 1.62, "Charlson", ha="center", va="center", fontsize=7.5, color="#64748B", zorder=4
        )

    def edge(x1, y1, x2, y2):
        ax.annotate(
            "",
            xy=(x2, y2),
            xytext=(x1, y1),
            arrowprops=dict(arrowstyle="-|>", color="#CBD5E1", lw=1.2, mutation_scale=12),
        )

    # Nodes
    node_icd(0.8)
    node_sig(2.5, "HR", "82 bpm")
    node_note(4.5, "NOTE 1")
    node_sig(6.5, "SpO₂", "97%")
    node_note(8.8, "NOTE 2")

    # Edges (temporal forward)
    edge(0.8, 1.8, 2.3, 1.8)
    edge(2.5, 2.3, 4.0, 1.85)
    edge(2.5, 1.8, 3.9, 1.55)
    edge(4.5, 1.8, 6.2, 1.8)
    edge(4.5, 1.2, 6.2, 1.5)
    edge(6.5, 1.8, 8.2, 1.65)
    edge(6.5, 2.3, 8.2, 1.9)

    # Edge attr label
    ax.text(3.5, 2.3, "Δt/24", fontsize=8, color="#94A3B8", ha="center", rotation=15)

    # Legend
    legend_items = [
        (plt.Circle((0, 0), 0.1, color="#06B6D4", alpha=0.15), "Sygnał (vital/lab) — 16-dim"),
        (
            FancyPatch.FancyBboxPatch(
                (0, 0), 0.2, 0.2, boxstyle="round", facecolor="#EEF2FF", edgecolor="#818CF8"
            ),
            "Notatka — 128-dim",
        ),
        (
            plt.Polygon(
                [(0, 0.1), (0.1, 0), (0.2, 0.1), (0.1, 0.2)],
                facecolor="#EEF2FF",
                edgecolor="#818CF8",
            ),
            "ICD Charlson — 19-dim",
        ),
    ]
    legend_proxies = [
        mpatches.Patch(color="#06B6D4", alpha=0.5, label="Sygnał (vital/lab) — 16-dim"),
        mpatches.Patch(color="#818CF8", alpha=0.5, label="Notatka kliniczna — 128-dim"),
        mpatches.Patch(color="#6366F1", alpha=0.3, label="ICD Charlson — 19-dim"),
    ]
    ax.legend(
        handles=legend_proxies,
        loc="upper right",
        fontsize=9,
        facecolor="white",
        edgecolor="#E2E8F0",
        labelcolor="#475569",
    )

    ax.text(
        0.1,
        3.65,
        "PHASE 2 — Budowanie grafu pacjenta (per pobyt OIOM)",
        fontsize=12,
        fontweight="bold",
        color="#0F172A",
    )

    fig.tight_layout(pad=0.1)
    return _png(fig)


def diagram_gnn() -> io.BytesIO:
    fig, ax = plt.subplots(figsize=(10, 3.8))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 3.8)
    ax.axis("off")

    def box(x, y, w, h, fc, ec, text, sub=None, tc="#0F172A"):
        rect = FancyPatch.FancyBboxPatch(
            (x, y),
            w,
            h,
            boxstyle="round,pad=0.06",
            facecolor=fc,
            edgecolor=ec,
            linewidth=1.5,
            zorder=3,
        )
        ax.add_patch(rect)
        ty = y + h - 0.28 if sub else y + h / 2
        ax.text(
            x + w / 2,
            ty,
            text,
            ha="center",
            va="top",
            fontsize=10,
            fontweight="bold",
            color=tc,
            zorder=4,
        )
        if sub:
            ax.text(
                x + w / 2,
                y + h / 2 - 0.12,
                sub,
                ha="center",
                va="center",
                fontsize=8.5,
                color="#64748B",
                zorder=4,
                linespacing=1.5,
            )

    def arrow(x1, x2, y=1.9):
        ax.annotate(
            "",
            xy=(x2, y),
            xytext=(x1, y),
            arrowprops=dict(arrowstyle="-|>", color="#CBD5E1", lw=1.5, mutation_scale=12),
        )

    # Flow boxes
    box(0.1, 1.1, 1.7, 1.7, "#F8FAFC", "#E2E8F0", "Input\nnody", "sig→16d\nnote→128d\nicd→19d")
    arrow(1.8, 2.15)
    box(
        2.2,
        0.8,
        1.8,
        2.1,
        "#EEF2FF",
        "#C7D2FE",
        "GINEConv",
        "×3 layers\n64→128→128\n+ReLU+Drop",
        "#4338CA",
    )
    arrow(4.0, 4.35)
    box(
        4.4,
        1.1,
        1.8,
        1.7,
        "#EFF6FF",
        "#BFDBFE",
        "Pooling",
        "attention\nor dual\nor mean",
        "#1E40AF",
    )
    arrow(6.2, 6.55)
    box(6.6, 1.4, 1.3, 1.1, "#EFF6FF", "#BFDBFE", "⊕ demo", "4-dim\n[age,sex\ntriage]", "#1E3A8A")
    arrow(7.9, 8.25)
    box(
        8.3,
        1.2,
        1.6,
        1.4,
        "#F1F5F9",
        "#CBD5E1",
        "Classifier",
        "132→32→1\nFocal Loss\nAUROC 0.850",
        "#475569",
    )

    # GINEConv formula
    ax.text(
        3.1,
        0.35,
        "h_v' = MLP(h_v + Σ(h_u + edge_attr))",
        ha="center",
        fontsize=9,
        color="#6366F1",
        style="italic",
    )

    # Param count
    ax.text(
        5.0, 3.6, "~64 700 parametrów", ha="center", fontsize=10, fontweight="bold", color="#6366F1"
    )

    ax.text(0.1, 3.7, "PHASE 2 — Architektura GNN", fontsize=12, fontweight="bold", color="#0F172A")

    fig.tight_layout(pad=0.1)
    return _png(fig)


# ─── PPTX helpers ─────────────────────────────────────────────────────────────


def new_slide(prs: Presentation) -> any:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = C_BG
    return slide


def txt(
    slide,
    text: str,
    x,
    y,
    w,
    h,
    size=20,
    bold=False,
    color=None,
    align=PP_ALIGN.LEFT,
    italic=False,
    wrap=True,
    spacing=1.0,
):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.line_spacing = spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = FONT
    run.font.color.rgb = color if color else C_TEXT
    return txb


def line(slide, x1, y1, x2, y2, color=None, width=Pt(0.75)):
    shape = slide.shapes.add_connector(1, x1, y1, x2, y2)
    shape.line.color.rgb = color if color else C_RULE
    shape.line.width = width


def rect_card(slide, x, y, w, h, fill=None, border=None, radius=Pt(6)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill if fill else C_CARD
    shape.line.color.rgb = border if border else C_BORDER
    shape.line.width = Pt(0.75)
    return shape


def img(slide, buf: io.BytesIO, x, y, w, h=None):
    buf.seek(0)
    if h:
        slide.shapes.add_picture(buf, x, y, w, h)
    else:
        slide.shapes.add_picture(buf, x, y, w)


def tag(slide, text: str, x, y):
    """Small pill tag."""
    w, h = Inches(1.5), Inches(0.28)
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb("EEECFB")
    shape.line.color.rgb = rgb("C7D2FE")
    shape.line.width = Pt(0.5)
    txf = shape.text_frame
    txf.word_wrap = False
    p = txf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text.upper()
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.name = FONT
    run.font.color.rgb = C_ACCENT


def stat_box(slide, number: str, label: str, x, y, w=Inches(2.2), accent=None):
    h = Inches(1.3)
    rect_card(slide, x, y, w, h)
    txt(
        slide,
        number,
        x + Inches(0.15),
        y + Inches(0.08),
        w - Inches(0.3),
        Inches(0.75),
        size=32,
        bold=True,
        color=accent if accent else C_INK,
        align=PP_ALIGN.CENTER,
    )
    txt(
        slide,
        label,
        x + Inches(0.1),
        y + Inches(0.82),
        w - Inches(0.2),
        Inches(0.38),
        size=11,
        color=C_MUTED,
        align=PP_ALIGN.CENTER,
    )


def bullet_row(slide, text: str, x, y, w, accent_dot=True):
    if accent_dot:
        txt(slide, "→", x, y, Inches(0.3), Inches(0.4), size=14, bold=True, color=C_ACCENT)
    txt(slide, text, x + Inches(0.3), y, w - Inches(0.3), Inches(0.45), size=14, color=C_TEXT)


# ─── Slides ───────────────────────────────────────────────────────────────────

PAD = Inches(0.7)
PAD2 = Inches(1.0)


def slide_title(prs, charts):
    s = new_slide(prs)
    CW = Inches(13.33)  # slide width

    # Thin accent bar at top (like NeuroNauts)
    top_bar = s.shapes.add_shape(1, Inches(0.55), Inches(0.38), Inches(0.4), Inches(0.05))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = C_ACCENT
    top_bar.line.fill.background()

    # Course tag centred
    txt(
        s,
        "GGSN  ·  2024/25",
        Inches(0),
        Inches(0.3),
        CW,
        Inches(0.3),
        size=11,
        bold=True,
        color=C_ACCENT,
        align=PP_ALIGN.CENTER,
    )

    # Large title — centred
    txt(
        s,
        "Temporal GNN",
        Inches(0),
        Inches(1.1),
        CW,
        Inches(1.25),
        size=60,
        bold=True,
        color=C_INK,
        align=PP_ALIGN.CENTER,
    )
    txt(
        s,
        "dla predykcji śmiertelności szpitalnej",
        Inches(0),
        Inches(2.3),
        CW,
        Inches(0.7),
        size=28,
        bold=False,
        color=C_MUTED,
        align=PP_ALIGN.CENTER,
    )

    # Horizontal rule
    line(s, Inches(2.0), Inches(3.15), Inches(11.33), Inches(3.15))

    # Subtitle centred
    txt(
        s,
        "MIMIC-IV 3.1  ·  Bio_ClinicalBERT  ·  GINEConv  ·  Phase 1 + Phase 2",
        Inches(0),
        Inches(3.25),
        CW,
        Inches(0.4),
        size=13,
        color=C_LIGHT,
        align=PP_ALIGN.CENTER,
    )

    # 3 stat boxes centred (no AUROC)
    # 3 × 3.3" + 2 × 0.5" gap = 10.9" → start at (13.33-10.9)/2 = 1.215"
    sx = Inches(1.215)
    bw = Inches(3.3)
    gap = Inches(0.5)
    stat_box(s, "52 727", "pobytów OIOM", sx, Inches(3.9), w=bw)
    stat_box(s, "12.7%", "śmiertelność\n(1 : 6.9 imbalans)", sx + bw + gap, Inches(3.9), w=bw)
    stat_box(s, "13", "ablacji", sx + 2 * (bw + gap), Inches(3.9), w=bw, accent=C_ACCENT)

    # Bottom rule + footer
    line(s, Inches(0.55), Inches(6.55), Inches(12.78), Inches(6.55))
    txt(
        s,
        "Temporal GNN  ·  MIMIC-IV 3.1",
        Inches(0.55),
        Inches(6.65),
        Inches(6),
        Inches(0.3),
        size=10,
        color=C_LIGHT,
    )
    txt(
        s,
        "GGSN 2024/25",
        Inches(7),
        Inches(6.65),
        Inches(5.78),
        Inches(0.3),
        size=10,
        color=C_LIGHT,
        align=PP_ALIGN.RIGHT,
    )


def slide_problem(prs, charts):
    s = new_slide(prs)
    tag(s, "Problem", PAD, Inches(0.5))
    txt(
        s,
        "Trzy główne wyzwania",
        PAD,
        Inches(0.85),
        Inches(7),
        Inches(0.7),
        size=34,
        bold=True,
        color=C_INK,
    )
    line(s, PAD, Inches(1.55), Inches(12.63), Inches(1.55))

    bullets = [
        ("Dane heterogeniczne", "notatki lekarskie + sygnały fizjologiczne + kody ICD-10"),
        ("Nierównoważne klasy", "1 : 6.9 — 5 343 zgony / 36 853 przeżycia"),
        ("Temporal leakage", "latwo nieumyslnie 'widziec przyszlosc' w danych sekwencyjnych"),
    ]
    for i, (title, sub) in enumerate(bullets):
        y = Inches(1.75 + i * 1.35)
        rect_card(s, PAD, y, Inches(7.5), Inches(1.1))
        txt(
            s,
            title,
            PAD + Inches(0.2),
            y + Inches(0.1),
            Inches(7),
            Inches(0.45),
            size=16,
            bold=True,
            color=C_INK,
        )
        txt(
            s,
            sub,
            PAD + Inches(0.2),
            y + Inches(0.52),
            Inches(7),
            Inches(0.45),
            size=13,
            color=C_MUTED,
        )

    # Side stats
    stat_box(
        s,
        "1 : 6.9",
        "imbalans klas\npos_weight = 6.897",
        Inches(8.8),
        Inches(1.75),
        w=Inches(4.0),
        accent=rgb("EF4444"),
    )
    stat_box(
        s,
        "0.127",
        "AUPRC random\nbaseline",
        Inches(8.8),
        Inches(3.25),
        w=Inches(4.0),
        accent=C_MUTED,
    )
    stat_box(
        s,
        "0.500",
        "AUROC random\nbaseline",
        Inches(8.8),
        Inches(4.75),
        w=Inches(4.0),
        accent=C_MUTED,
    )


def slide_data(prs, charts):
    s = new_slide(prs)
    tag(s, "Dane", PAD, Inches(0.5))
    txt(
        s,
        "MIMIC-IV 3.1 — trzy modalności",
        PAD,
        Inches(0.85),
        Inches(11),
        Inches(0.7),
        size=34,
        bold=True,
        color=C_INK,
    )
    line(s, PAD, Inches(1.55), Inches(12.63), Inches(1.55))

    cols = [
        ("Notatki kliniczne", "102 221", "unikalnych notatek\nBio_ClinicalBERT tokenizacja"),
        ("Sygnały fizjologiczne", "55.7M", "wierszy (vitals + labs)\n14 typów: HR, SpO₂, MAP, ..."),
        ("Diagnozy ICD-10", "19", "kategorii Charlson\nbinarny wektor per pobyt"),
    ]
    for i, (title, number, sub) in enumerate(cols):
        x = PAD + Inches(i * 4.0)
        rect_card(s, x, Inches(1.85), Inches(3.7), Inches(4.5))
        txt(
            s,
            title,
            x + Inches(0.2),
            Inches(2.1),
            Inches(3.4),
            Inches(0.5),
            size=14,
            bold=True,
            color=C_INK,
        )
        txt(
            s,
            number,
            x + Inches(0.2),
            Inches(2.8),
            Inches(3.4),
            Inches(1.2),
            size=44,
            bold=True,
            color=C_ACCENT,
            align=PP_ALIGN.LEFT,
        )
        txt(s, sub, x + Inches(0.2), Inches(4.2), Inches(3.4), Inches(1.0), size=12, color=C_MUTED)

    # Footer note
    rect_card(
        s, PAD, Inches(6.5), Inches(11.93), Inches(0.7), fill=rgb("EEF2FF"), border=rgb("C7D2FE")
    )
    txt(
        s,
        "Podział subject-disjoint — pacjent nigdy jednocześnie w zbiorze treningowym i walidacyjnym",
        PAD + Inches(0.2),
        Inches(6.6),
        Inches(11.5),
        Inches(0.5),
        size=13,
        color=C_ACCENT,
    )


def slide_architecture(prs, charts):
    s = new_slide(prs)
    tag(s, "Architektura", PAD, Inches(0.5))
    txt(
        s,
        "Dwufazowe podejście",
        PAD,
        Inches(0.85),
        Inches(9),
        Inches(0.7),
        size=34,
        bold=True,
        color=C_INK,
    )
    line(s, PAD, Inches(1.55), Inches(12.63), Inches(1.55))

    phases = [
        (
            "PHASE 1",
            "Contrastive Pre-training",
            "Two-Tower (TextTower + SignalTower)\nInfoNCE · τ=0.07 · grad accum (eff. batch 64)\n→ 102K zamrożonych embeddingów 128-dim",
            rgb("EEF2FF"),
            rgb("C7D2FE"),
            C_ACCENT,
        ),
        (
            "PHASE 2",
            "Temporal Heterogeneous GNN",
            "Graf pacjenta (signal + note + ICD nody)\nGINEConv×3 · attention pooling\nFocal Loss · demographics",
            rgb("ECFDF5"),
            rgb("A7F3D0"),
            C_GREEN,
        ),
    ]
    for i, (phase, title, body, fc, ec, tc) in enumerate(phases):
        x = PAD + Inches(i * 5.8)
        rect_card(s, x, Inches(1.85), Inches(5.5), Inches(4.5), fill=fc, border=ec)
        txt(
            s,
            phase,
            x + Inches(0.2),
            Inches(2.0),
            Inches(5.0),
            Inches(0.4),
            size=11,
            bold=True,
            color=tc,
        )
        txt(
            s,
            title,
            x + Inches(0.2),
            Inches(2.4),
            Inches(5.0),
            Inches(0.6),
            size=19,
            bold=True,
            color=C_INK,
        )
        txt(
            s,
            body,
            x + Inches(0.2),
            Inches(3.0),
            Inches(5.0),
            Inches(1.8),
            size=13,
            color=C_MUTED,
            spacing=1.4,
        )

    # Arrow between phases
    txt(
        s,
        "→",
        Inches(6.35),
        Inches(3.7),
        Inches(0.7),
        Inches(0.6),
        size=36,
        bold=True,
        color=C_LIGHT,
        align=PP_ALIGN.CENTER,
    )

    # Result
    rect_card(
        s, PAD, Inches(6.5), Inches(11.93), Inches(0.7), fill=rgb("FFFBEB"), border=rgb("FDE68A")
    )
    txt(
        s,
        "Wynik: AUROC 0.850 · AUPRC 0.465 · ~64 700 parametrów",
        PAD + Inches(0.2),
        Inches(6.6),
        Inches(11.5),
        Inches(0.5),
        size=14,
        bold=True,
        color=rgb("92400E"),
    )


def slide_phase1(prs, charts):
    s = new_slide(prs)
    tag(s, "Phase 1", PAD, Inches(0.5))
    txt(
        s,
        "Contrastive Pre-training",
        PAD,
        Inches(0.85),
        Inches(11),
        Inches(0.7),
        size=34,
        bold=True,
        color=C_INK,
    )
    line(s, PAD, Inches(1.55), Inches(12.63), Inches(1.55))
    img(s, charts["twotower"], PAD, Inches(1.8), Inches(8.5))

    # Right column — kluczowe liczby Phase 1
    x2 = PAD + Inches(8.9)
    stat_box(
        s,
        "2.980",
        "best val InfoNCE\n(−28.3% vs baseline)",
        x2,
        Inches(1.8),
        w=Inches(3.7),
        accent=C_ACCENT,
    )
    stat_box(s, "0.730", "linear probe AUROC", x2, Inches(3.3), w=Inches(3.7), accent=C_GREEN)
    stat_box(
        s,
        "102 221",
        "zamrożonych embeddingów\n(128-dim)",
        x2,
        Inches(4.8),
        w=Inches(3.7),
        accent=C_CYAN,
    )

    txt(
        s,
        "Trening: Bio_ClinicalBERT (top 4/12 fine-tunowane) · 16 epok · early stop",
        PAD,
        Inches(6.85),
        Inches(11.93),
        Inches(0.35),
        size=11,
        italic=True,
        color=C_LIGHT,
        align=PP_ALIGN.CENTER,
    )


def slide_phase2_graph(prs, charts):
    s = new_slide(prs)
    tag(s, "Phase 2 — Graf", PAD, Inches(0.5))
    txt(
        s,
        "Budowanie grafu pacjenta",
        PAD,
        Inches(0.85),
        Inches(11),
        Inches(0.7),
        size=34,
        bold=True,
        color=C_INK,
    )
    line(s, PAD, Inches(1.55), Inches(12.63), Inches(1.55))
    img(s, charts["graph"], PAD, Inches(1.7), Inches(12.0))

    # Small annotations
    bullets_text = [
        "Krawędzie skierowane tylko do przodu w czasie",
        "Atrybut krawędzi: Δt/24 (różnica w dobach)",
        "ICD node (t=−1h) ← wiedza a priori",
    ]
    for i, b in enumerate(bullets_text):
        bullet_row(s, b, PAD, Inches(5.95 + i * 0.45), Inches(11.9), accent_dot=True)


def slide_gnn(prs, charts):
    s = new_slide(prs)
    tag(s, "Phase 2 — GNN", PAD, Inches(0.5))
    txt(
        s,
        "Architektura sieci grafowej",
        PAD,
        Inches(0.85),
        Inches(11),
        Inches(0.7),
        size=34,
        bold=True,
        color=C_INK,
    )
    line(s, PAD, Inches(1.55), Inches(12.63), Inches(1.55))
    img(s, charts["gnn"], PAD, Inches(1.7), Inches(12.0))
    txt(
        s,
        "Implementacja: PyTorch Lightning + PyTorch Geometric",
        PAD,
        Inches(6.85),
        Inches(11.93),
        Inches(0.35),
        size=11,
        italic=True,
        color=C_LIGHT,
        align=PP_ALIGN.CENTER,
    )


def slide_bug(prs, charts):
    s = new_slide(prs)
    tag(s, "Wykryty błąd", PAD, Inches(0.5))
    txt(
        s,
        "Temporal Leakage — AUROC 0.931 = alarm",
        PAD,
        Inches(0.85),
        Inches(11),
        Inches(0.8),
        size=30,
        bold=True,
        color=C_INK,
    )
    line(s, PAD, Inches(1.65), Inches(12.63), Inches(1.65))

    img(s, charts["leakage"], PAD, Inches(1.8), Inches(5.8))

    # Right side cards
    x2 = Inches(7.1)
    rect_card(
        s, x2, Inches(1.8), Inches(5.5), Inches(2.4), fill=rgb("FEF2F2"), border=rgb("FECACA")
    )
    txt(
        s,
        "Stary kod",
        x2 + Inches(0.2),
        Inches(1.95),
        Inches(5),
        Inches(0.4),
        size=14,
        bold=True,
        color=C_RED,
    )
    txt(
        s,
        "sig_rows[-max_signals:]",
        x2 + Inches(0.2),
        Inches(2.4),
        Inches(5),
        Inches(0.35),
        size=12,
        italic=True,
        color=C_MUTED,
    )
    txt(
        s,
        "= ostatnie N sygnałów = sygnały tuż przed śmiercią\nAUROC 0.931 — trivially easy task",
        x2 + Inches(0.2),
        Inches(2.8),
        Inches(5),
        Inches(0.8),
        size=12,
        color=C_TEXT,
    )

    rect_card(
        s, x2, Inches(4.4), Inches(5.5), Inches(2.4), fill=rgb("F0FDF4"), border=rgb("A7F3D0")
    )
    txt(
        s,
        "Poprawka",
        x2 + Inches(0.2),
        Inches(4.55),
        Inches(5),
        Inches(0.4),
        size=14,
        bold=True,
        color=C_GREEN,
    )
    txt(
        s,
        "sig_rows[:max_signals]",
        x2 + Inches(0.2),
        Inches(5.0),
        Inches(5),
        Inches(0.35),
        size=12,
        italic=True,
        color=C_MUTED,
    )
    txt(
        s,
        "= pierwsze N sygnałów = wczesna faza pobytu\nAUROC 0.820 — realny wynik",
        x2 + Inches(0.2),
        Inches(5.4),
        Inches(5),
        Inches(0.8),
        size=12,
        color=C_TEXT,
    )


def slide_auroc(prs, charts):
    s = new_slide(prs)
    tag(s, "Wyniki", PAD, Inches(0.5))
    txt(
        s,
        "AUROC — 13 eksperymentów",
        PAD,
        Inches(0.85),
        Inches(11),
        Inches(0.65),
        size=34,
        bold=True,
        color=C_INK,
    )
    line(s, PAD, Inches(1.5), Inches(12.63), Inches(1.5))
    img(s, charts["auroc"], Inches((13.33 - 11.0) / 2), Inches(1.65), Inches(11.0))


def slide_auprc(prs, charts):
    s = new_slide(prs)
    tag(s, "Wyniki", PAD, Inches(0.5))
    txt(
        s,
        "AUPRC — 3.6× powyżej losowego",
        PAD,
        Inches(0.85),
        Inches(11),
        Inches(0.65),
        size=34,
        bold=True,
        color=C_INK,
    )
    line(s, PAD, Inches(1.5), Inches(12.63), Inches(1.5))
    img(s, charts["auprc"], Inches((13.33 - 11.0) / 2), Inches(1.65), Inches(11.0))

    txt(
        s,
        "AUPRC ważniejsza niż AUROC przy silnym imbalansie klas (1:7)",
        PAD,
        Inches(7.0),
        Inches(11.93),
        Inches(0.35),
        size=12,
        color=C_LIGHT,
        align=PP_ALIGN.CENTER,
    )


def slide_contributions(prs, charts):
    s = new_slide(prs)
    tag(s, "Analiza", PAD, Inches(0.5))
    txt(
        s,
        "Wkład każdego komponentu",
        PAD,
        Inches(0.85),
        Inches(8),
        Inches(0.65),
        size=34,
        bold=True,
        color=C_INK,
    )
    line(s, PAD, Inches(1.5), Inches(12.63), Inches(1.5))

    img(s, charts["contributions"], PAD, Inches(1.7), Inches(7.0))

    # Right side
    x2 = Inches(7.95)
    rect_card(
        s, x2, Inches(1.7), Inches(4.7), Inches(2.5), fill=rgb("F0FDF4"), border=rgb("A7F3D0")
    )
    txt(
        s,
        "Co działa",
        x2 + Inches(0.2),
        Inches(1.85),
        Inches(4.3),
        Inches(0.4),
        size=13,
        bold=True,
        color=C_GREEN,
    )
    works = [
        "Notatki kliniczne:  +0.045",
        "Demografika:  +0.012",
        "Attention pooling:  +0.015",
        "Focal loss:  lepsza sens@95spec",
    ]
    for i, w in enumerate(works):
        txt(
            s,
            w,
            x2 + Inches(0.3),
            Inches(2.3 + i * 0.4),
            Inches(4.3),
            Inches(0.38),
            size=12,
            color=C_TEXT,
        )

    rect_card(
        s, x2, Inches(4.4), Inches(4.7), Inches(2.5), fill=rgb("FEF2F2"), border=rgb("FECACA")
    )
    txt(
        s,
        "Co nie działa",
        x2 + Inches(0.2),
        Inches(4.55),
        Inches(4.3),
        Inches(0.4),
        size=13,
        bold=True,
        color=C_RED,
    )
    fails = [
        "Hard negatives: catastrophic forgetting",
        "ICD node: 56% stays = zero",
        "All-stay signals: nie bije ±2h",
        "E2E BERT: catastrophic forgetting",
    ]
    for i, f in enumerate(fails):
        txt(
            s,
            f,
            x2 + Inches(0.3),
            Inches(5.0 + i * 0.4),
            Inches(4.3),
            Inches(0.38),
            size=12,
            color=C_TEXT,
        )


def slide_best_model(prs, charts):
    s = new_slide(prs)
    tag(s, "Najlepszy model", PAD, Inches(0.5))
    txt(
        s,
        "demo_attention_focal2",
        PAD,
        Inches(0.85),
        Inches(9),
        Inches(0.7),
        size=36,
        bold=True,
        color=C_INK,
    )
    line(s, PAD, Inches(1.55), Inches(12.63), Inches(1.55))

    stat_box(s, "0.850", "AUROC", PAD, Inches(1.8), w=Inches(3.5), accent=C_ACCENT)
    stat_box(s, "0.465", "AUPRC", PAD + Inches(3.7), Inches(1.8), w=Inches(3.5), accent=C_GREEN)
    stat_box(
        s, "0.388", "sens @ 95% spec", PAD + Inches(7.4), Inches(1.8), w=Inches(4.9), accent=C_CYAN
    )

    # Config card — bullet list
    rect_card(s, PAD, Inches(3.3), Inches(11.93), Inches(2.7), fill=rgb("F8FAFC"), border=C_RULE)
    txt(
        s,
        "Konfiguracja",
        PAD + Inches(0.2),
        Inches(3.4),
        Inches(11.5),
        Inches(0.4),
        size=13,
        color=C_ACCENT,
        bold=True,
    )

    config_rows = [
        "Pooling: attention",
        "Focal loss γ = 2.0",
        "Demographics (4-dim, graph-level po poolingu)",
        "hidden_dim = 128 · 3× GINEConv · dropout 0.3",
        "Optymizator: Adam · lr = 1e-3 · ReduceLROnPlateau",
    ]
    for i, row in enumerate(config_rows):
        bullet_row(s, row, PAD + Inches(0.2), Inches(3.85 + i * 0.4), Inches(11.5), accent_dot=True)

    stat_box(
        s, "3.6×", "AUPRC powyżej\nlosowego", PAD, Inches(6.15), w=Inches(3.5), accent=C_ACCENT
    )
    stat_box(
        s,
        "1:7",
        "imbalans klas\n(pos_weight=6.9)",
        PAD + Inches(3.7),
        Inches(6.15),
        w=Inches(3.5),
        accent=rgb("EF4444"),
    )
    stat_box(
        s,
        "~65K",
        "parametrów\n(vs 110M BERT)",
        PAD + Inches(7.4),
        Inches(6.15),
        w=Inches(4.9),
        accent=C_MUTED,
    )


def slide_sota(prs, charts):
    s = new_slide(prs)
    tag(s, "Kontekst", PAD, Inches(0.5))
    txt(
        s,
        "Nasz wynik vs SOTA",
        PAD,
        Inches(0.85),
        Inches(9),
        Inches(0.7),
        size=34,
        bold=True,
        color=C_INK,
    )
    line(s, PAD, Inches(1.55), Inches(12.63), Inches(1.55))

    img(s, charts["sota"], PAD, Inches(1.7), Inches(5.8))

    x2 = Inches(7.1)
    rect_card(
        s, x2, Inches(1.8), Inches(5.5), Inches(2.2), fill=rgb("FEF2F2"), border=rgb("FECACA")
    )
    txt(
        s,
        "Dlaczego gap 0.030?",
        x2 + Inches(0.2),
        Inches(1.95),
        Inches(5),
        Inches(0.4),
        size=13,
        bold=True,
        color=C_RED,
    )
    for i, b in enumerate(
        [
            "SOTA: end-to-end BERT (110M params, multi-GPU)",
            "My: zamrożone embeddingi → bottleneck",
            "E2E próbowane → catastrophic forgetting",
        ]
    ):
        txt(
            s,
            b,
            x2 + Inches(0.3),
            Inches(2.4 + i * 0.45),
            Inches(5.0),
            Inches(0.4),
            size=12,
            color=C_TEXT,
        )

    rect_card(
        s, x2, Inches(4.2), Inches(5.5), Inches(2.5), fill=rgb("F0FDF4"), border=rgb("A7F3D0")
    )
    txt(
        s,
        "Dlaczego to dobry wynik?",
        x2 + Inches(0.2),
        Inches(4.35),
        Inches(5),
        Inches(0.4),
        size=13,
        bold=True,
        color=C_GREEN,
    )
    for i, b in enumerate(
        [
            "~65K params vs 110M BERT — ~1 700× mniejszy",
            "AUPRC 0.465 = 3.6× powyżej losowego",
            "13 ablacji — wkład każdego komponentu zmierzony",
            "Negatywne wyniki też opisane (E2E, ICD, hard negs)",
        ]
    ):
        txt(
            s,
            b,
            x2 + Inches(0.3),
            Inches(4.85 + i * 0.41),
            Inches(5.0),
            Inches(0.38),
            size=12,
            color=C_TEXT,
        )


def slide_takeaways(prs, charts):
    s = new_slide(prs)
    tag(s, "Wnioski", PAD, Inches(0.5))
    txt(
        s,
        "5 kluczowych lekcji",
        PAD,
        Inches(0.85),
        Inches(11),
        Inches(0.65),
        size=34,
        bold=True,
        color=C_INK,
    )
    line(s, PAD, Inches(1.5), Inches(12.63), Inches(1.5))

    takeaways = [
        (
            "01",
            C_ACCENT,
            "Multimodalne GNN działają",
            "Ablacja signal_only: notatki kliniczne dają +0.045 AUROC — obie modalności są kluczowe",
        ),
        (
            "02",
            C_GREEN,
            "Phase 1 pre-training → dobre zamrożone reprezentacje",
            "E2E fine-tuning niszczy reprezentacje (catastrophic forgetting, AUROC 0.790 < 0.829)",
        ),
        (
            "03",
            rgb("F59E0B"),
            "Temporal leakage jest subtelny i łatwy do przeoczenia",
            "Ostatnie N sygnałów = sygnały przed śmiercią → AUROC 0.931 (artefakt). Pierwsze N → 0.820",
        ),
        (
            "04",
            C_CYAN,
            "Demografika to prosty ale skuteczny bonus",
            "4 cechy [wiek, płeć, tryb przyjęcia] dają +0.012 AUROC jako graph-level feature po poolingu",
        ),
        (
            "05",
            C_RED,
            "Negatywne wyniki potwierdziły dobór finalnego stacku",
            "Hard negatives, ICD node, all-stay signals, E2E — każde gorsze od baseline",
        ),
    ]

    for i, (num, nc, title, body) in enumerate(takeaways):
        y = Inches(1.65 + i * 1.02)
        rect_card(s, PAD, y, Inches(11.93), Inches(0.88))
        txt(
            s,
            num,
            PAD + Inches(0.15),
            y + Inches(0.08),
            Inches(0.6),
            Inches(0.7),
            size=22,
            bold=True,
            color=nc,
        )
        txt(
            s,
            title,
            PAD + Inches(0.8),
            y + Inches(0.06),
            Inches(11),
            Inches(0.38),
            size=15,
            bold=True,
            color=C_INK,
        )
        txt(
            s,
            body,
            PAD + Inches(0.8),
            y + Inches(0.43),
            Inches(11),
            Inches(0.38),
            size=12,
            color=C_MUTED,
        )


def slide_umap(prs, charts):
    s = new_slide(prs)
    tag(s, "Phase 1 — walidacja jakości", PAD, Inches(0.5))
    txt(
        s,
        "Embeddingi separują pacjentów według outcome",
        PAD,
        Inches(0.85),
        Inches(11),
        Inches(0.7),
        size=28,
        bold=True,
        color=C_INK,
    )
    line(s, PAD, Inches(1.55), Inches(12.63), Inches(1.55))

    img(s, charts["umap"], PAD, Inches(1.7), Inches(8.5))

    x2 = PAD + Inches(8.9)
    rect_card(s, x2, Inches(1.8), Inches(3.7), Inches(4.8), fill=rgb("F8FAFC"), border=C_RULE)
    txt(
        s,
        "Co tu widać",
        x2 + Inches(0.2),
        Inches(1.95),
        Inches(3.4),
        Inches(0.4),
        size=13,
        bold=True,
        color=C_ACCENT,
    )

    bullets = [
        "UMAP 2D zamrożonych embeddingów (128-dim)",
        "Zgon (czerwony) klastruje w obszarach o wyższym ryzyku — sygnał obecny już w pre-trainingu",
        "Bez fine-tuningu na mortality — pure InfoNCE pre-training",
        "Wyjaśnia +0.045 AUROC vs signal_only baseline",
    ]
    for i, b in enumerate(bullets):
        bullet_row(s, b, x2 + Inches(0.2), Inches(2.4 + i * 1.0), Inches(3.4), accent_dot=True)


def slide_learning_curve(prs, charts):
    s = new_slide(prs)
    tag(s, "Walidacja modelu", PAD, Inches(0.5))
    txt(
        s,
        "Krzywa uczenia — brak overfittingu",
        PAD,
        Inches(0.85),
        Inches(11),
        Inches(0.7),
        size=30,
        bold=True,
        color=C_INK,
    )
    line(s, PAD, Inches(1.55), Inches(12.63), Inches(1.55))

    img(s, charts["learning_curve"], Inches((13.33 - 11.0) / 2), Inches(1.7), Inches(11.0))

    txt(
        s,
        "val AUROC plateauje ~0.845, train loss dalej maleje — model uczy się "
        "bez nadmiernego dopasowania. Early stopping (patience=15) potwierdza stabilność.",
        PAD,
        Inches(6.9),
        Inches(11.93),
        Inches(0.4),
        size=12,
        italic=True,
        color=C_LIGHT,
        align=PP_ALIGN.CENTER,
    )


# ─── Main ─────────────────────────────────────────────────────────────────────


def main():
    print("Generowanie wykresów (light theme)…")
    charts = {
        "auroc": chart_auroc(),
        "auprc": chart_auprc(),
        "leakage": chart_leakage(),
        "contributions": chart_contributions(),
        "sota": chart_sota(),
        "learning_curve": chart_learning_curve(),
    }
    print("Generowanie diagramów…")
    charts["twotower"] = diagram_twotower()
    charts["graph"] = diagram_graph()
    charts["gnn"] = diagram_gnn()
    print(f"  {len(charts)} grafik gotowych.")

    print("Budowanie prezentacji PPTX…")
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slide_title(prs, charts)
    slide_problem(prs, charts)
    slide_data(prs, charts)
    slide_architecture(prs, charts)
    slide_phase1(prs, charts)
    slide_phase2_graph(prs, charts)
    slide_gnn(prs, charts)
    slide_bug(prs, charts)
    slide_auroc(prs, charts)
    slide_auprc(prs, charts)
    slide_learning_curve(prs, charts)
    slide_contributions(prs, charts)
    slide_best_model(prs, charts)
    slide_sota(prs, charts)
    slide_takeaways(prs, charts)

    # Notatki prelegenta
    assert len(prs.slides) == len(SLIDE_NOTES), (
        f"Liczba slajdów ({len(prs.slides)}) != liczba notatek ({len(SLIDE_NOTES)})"
    )
    for slide, notes in zip(prs.slides, SLIDE_NOTES):
        _add_notes(slide, notes)
    print(f"  Notatki: {len(SLIDE_NOTES)} slajdów.")

    prs.save(OUT)
    size_kb = OUT.stat().st_size // 1024
    print(f"\n  Zapisano: {OUT}  ({size_kb} KB)")
    print(f"  Slajdy: {len(prs.slides)}")


# ─── Speaker notes ────────────────────────────────────────────────────────────
# Kolejność odpowiada kolejności slajdów w main().

SLIDE_NOTES = [
    # 1 — Title
    """\
CO TO JEST MIMIC-IV?
Publiczna baza danych Beth Israel Deaconess Medical Center (Boston). Zawiera dane
EHR 94 429 unikalnych pacjentów. Dostep wymaga szkolenia z etyki przez PhysioNet.
Wersja 3.1 — najnowsza, uzyta w projekcie.

DLACZEGO PREDYKCJA SMIERTELNOSCI?
Kliniczne uzasadnienie: wczesne rozpoznanie pacjentow wysokiego ryzyka pozwala
na szybsze leczenie, priorytetyzacje zasobow OIOM, wlasciwy poziom opieki.
AUROC 0.88 SOTA (literatura) pochodzi z modeli z pelnymi zasobami obliczeniowymi;
nasz wynik 0.850 to realistyczny cel przy ograniczeniach jednego GPU laptopowego.

CZYM SIE ROZNI OD KLASYCZNEJ KLASYFIKACJI?
Dane sa sekwencyjne w czasie (zdarzenia), heterogeniczne (tekst + liczby + kody),
i maj silny imbalans klas. To wymaga specjalizowanej architektury.
""",
    # 2 — Problem
    """\
DLACZEGO 1:7 TO PROBLEM? NIE WYSTARCZY BCE?
BCE (binary cross-entropy) traktuje wszystkie przyklady rowno. Przy 1:7
model ktory zawsze mowi "przezyje" ma ~87% accuracy. Musimy uzyc:
1. pos_weight w losie — skaluje gradient dla klasy pozytywnej (pos_weight=6.897)
2. Focal Loss — dodaje czynnik (1-p)^gamma skupiajacy gradient na trudnych
   przykladach, niezaleznie od klasy

DLACZEGO AUROC A NIE ACCURACY?
AUROC mierzy zdolnosc rankingowania: czy model wyze rankuje przypadki smiertelne
powyzej ocalonych, niezaleznie od progu decyzyjnego. Jest odporna na imbalans.
AUPRC (precision-recall) jest jeszcze lepsza przy silnym imbalansie —
AUROC moze byc mylacy jesli wiele True Negatives zawyza metryki.

CO TO TEMPORAL LEAKAGE?
Sytuacja gdy model uzywa informacji ktora nie bylyby dostepna w czasie predykcji.
W naszym przypadku: wzielismy OSTATNIE N sygnalow z calego pobytu.
Pacjent ktory umiera ma krytyczne vitals bezposrednio przed smiercia —
model widzial "zapis umierania" zamiast prognozowac na podstawie wczesnych danych.
""",
    # 3 — Data
    """\
DLACZEGO BIO_CLINICALBERT A NIE ZWYKLY BERT?
Bio_ClinicalBERT (emilyalsentzer) byl dotrenowany na notatkach z MIMIC-III
i literaturze biomedycznej PubMed. Kluczowy: zna skroty kliniczne (PRN, QID,
SBP, MAP), terminologie medyczna, styl notatek pielegniarskich.
Ogolny BERT traktuje "WBC 12.5" jak zwykly ciag znakow; Bio_ClinicalBERT
rozumie ze to jest wynik morfologii.

DLACZEGO AKURAT 14 TYPOW SYGNALOW?
Wybrano 14 najczesciej wystepujacych i klinicznie istotnych sygnalow w MIMIC-IV:
HR (tetno), SpO2 (saturacja), MAP (cisnienie), RR (oddech), GCS (swiadomosc),
temperatura, i wybrane wyniki laboratoryjne. Mozna by wziac wiecej, ale 14
zapewnia pokrycie >99% pobytow przy zachowaniu rozsdnym rozmiarze modelu.

DLACZEGO SUBJECT-DISJOINT SPLIT A NIE LOSOWY?
Pacjent moze miec wiele pobytow na OIOM-ie. Losowy podzial moze dac ten sam
podmiot w train i val — model uczy sie cech specyficznych dla konkretnego
pacjenta (wiek, historia choroby), co sztucznie zawyza wyniki. Subject-disjoint
gwarantuje ze evalujemy generalizacje do nowych pacjentow, nie nowych pobytow.

CO TO CHARLSON COMORBIDITY INDEX?
Skala z 1987 roku (Charlson et al.), zmodyfikowana przez Quan et al. 2005 — 17
schorzeń przewleklych wazona wedlug ryzyka smiertelnosci (zawál, udar, nowotwory
przerzutowe...). Uzywamy wersji 19-kategoryjnej jako 0/1 wektora binarnego.
Wazne: kody ICD sa przypisywane PRZY WYPISIE, ale kategorie Charlson to
schorzenia istniejace przed hospitalizacja — brak leakage.
""",
    # 4 — Architecture
    """\
DLACZEGO DWIE FAZY A NIE END-TO-END OD RAZU?
E2E fine-tuning przetestowalismy (slajd o wynikach). Wynik: AUROC 0.790 < 0.829
baseline — catastrophic forgetting. BERT ma 110M parametrow, GNN ~65K.
Przy wspolnym treningu gradient z klasyfikatora "niszczy" reprezentacje BERT
wyuczone na 102K notatkach. Phase 1 daje stabilne, bogate reprezentacje
ktore Phase 2 moze efektywnie wykorzystac bez ich niszczenia.

CZYM ROZNI SIE OD FINE-TUNINGU BERT BEZPOSREDNIO?
Tradycyjny fine-tuning BERT na klasyfikacji: pobieramy [CLS] token i uczymy
klasyfikatora. Nie korzysta z sygnalow fizjologicznych ani kolejnosci zdarzen.
Nasze podejscie: Phase 1 tworzy WSPOLNA przestrzen tekst+sygnal, potem GNN
modeluje RELACJE TEMPORALNE miedzy zdarzeniami. To architektonicznie bogatsze.

JAKI JEST WKLAD KAZDEJ FAZY?
Phase 1: dostarcza reprezentacje notatek (bez niej: AUROC 0.784, delta -0.045)
Phase 2: modeluje kontekst temporalny i relacje miedzy zdarzeniami w czasie
Razem: AUROC 0.850
""",
    # 5 — Phase 1 Two-Tower
    """\
CO TO JEST INFONCE / NT-XENT?
InfoNCE (Noise Contrastive Estimation) minimalizuje temperatura-skalowane
cross-entropy na macierzy podobienstwa: L = CE(sim/tau, labels_diagonal).
Kazda para (notatka_i, sygnal_i) z tego samego pobytu jest para pozytywna,
wszystkie inne pary w batchu sa negatywne. Pozwala na O(B^2) negatywnych par
z B przykladow — efektywniejsze niz triplet loss.

DLACZEGO TEMPERATURA TAU=0.07?
Temperatura kontroluje "ostrosc" dystrybucji podobienstwa:
- Male tau (0.07): ostre piki, model musi byc bardzo pewny par pozytywnych
- Duze tau (>0.5): rozmyta dystrybucja, slabszy sygnal treningowy
Tau=0.07 pochodzi z SimCLR (Chen et al. 2020) i CLIP (Radford et al. 2021) —
empirycznie sprawdzone na multimodalnym uczeniu kontrastywnym.

DLACZEGO L2-NORMALIZACJA WEKTOROW?
L2-normalizacja projkuje wektory na sferze jednostkowej. Wtedy:
dot(z1, z2) = cosine_similarity(z1, z2) ∈ [-1, 1]
Zapobiega "collapse" (wszystkie wektory blisko zera), stabilizuje trening,
ulatwia interpretacje podobienstwa jako kosinusowe.

CO TO HARD NEGATIVE MINING?
W batchu szukamy par ktore sa trudne — tj. notatki i sygnaly z ROZNYCH pobytow
ale o wysokim podobienst. Takie "trudne negatywne" sa bardziej informacyjne
niz losowe (np. notatka z kardiologii vs sygnaly z neurologii — latwo rozroznic;
dwie notatki z podobnymi vitals z roznych pobytow — trudno). Hard negatywne
przyspieszaja uczenie sie wyraZnych granic w przestrzeni embeddingów.
""",
    # 6 — Phase 2 Graph
    """\
DLACZEGO GRAF A NIE TRANSFORMER / LSTM?
Transformer (attention) zaklada stala dlugosc sekwencji i rowno-rozmieszczone
punkty czasowe. Nasze zdarzenia sa nieregularnie rozmieszczone w czasie
(notatki co kilka godzin, sygnaly co minuty lub godziny).
Graf pozwala naturalnie reprezento wac nieregularne struktury czasowe.
LSTM przetwarza sekwencje liniowo — nie uwzglednia ze zdarzenia moga byc
z roznych typow (notatka vs sygnal) i w roznych momentach.

DLACZEGO KRAWEDZIE TYLKO DO PRZODU W CZASIE?
Krawedzie skierowane zapewniaja kauzalnosc: wczesniejsze zdarzenia moga
wplywac na interpretacje pozniejszych, ale nie odwrotnie. Cofanie krawedzi
byloby rownoważne z data leakage — model "widzialby przyszlosc" przez
propagacje wiadomosci wstecz.

DLACZEGO ICD NODE W T=-1H?
ICD Charlson reprezentuje PRZEWLEKLE schorzenia istniejace PRZED hospitalizacja.
Umieszczamy go przed wszystkimi zdarzeniami (t=-1h) jako "prior knowledge" —
wiedza a priori o stanie zdrowia pacjenta. Krawedz ICD→kazdy_node oznacza
ze historia chorobowa warunkuje interpretacje kazdego zdarzenia.

DLACZEGO DEMOGRAFIKA JAKO GRAPH-LEVEL A NIE NODE?
Wiek, plec, tryb przyjecia to cechy STATYCZNE dla calego pobytu — nie zdarzenia
w czasie. Dodanie ich jako osobnego nodu w grafie temporalnym byloby niespojne
architektonicznie. Jako graph-level feature (konkatenacja po poolingu) sa
dostepne klasyfikatorowi bez wplywania na propagacje wiadomosci w grafie.
""",
    # 7 — GNN Architecture
    """\
DLACZEGO GINECONV A NIE GAT (Graph Attention Network)?
GAT (Velickovic 2018) uczy sie wag uwagi dla krawedzi na podstawie cech NODOW.
Problem: nie uwzglednia natywnie cech krawedzi (u nas: Δt).
GINEConv (Hu et al. 2019) rozszerza GIN o cechy krawedzi: dodaje edge_attr do
cechy nodu przed agregacja. Poniewaz Δt (roznica czasu) jest kluczowa semantycznie
— odleglosc 30 minut vs 12 godzin miedzy notatka a sygnalami ma inne znaczenie —
potrzebujemy modelu ktory te rozne moze explicite przetworzyc.

DLACZEGO GIN A NIE GCN?
GCN (Kipf & Welling 2017) uzywca normalizacji stopni wezlow do agregacji.
Jest mniej ekspresywny — nie rozroznia grafow ktore GIN rozroznia.
GIN (Xu et al. 2019) jest tak ekspresywny jak test Weisfeiler-Leman (WL-test) —
teoretyczne maksimum dla message-passing GNN. MLP wewnatrz konwolucji zamiast
prostej sumy daje dodatkowa ekspresywnosc.

DLACZEGO 3 WARSTWY A NIE WIECEJ?
Ablacja pokazala: 2 warstwy → AUROC 0.835, 3 warstwy → 0.829 (baseline), 4 warstwy → 0.822.
Z demografika: demo+4layers → 0.836 (gorsze niz demo+3layers → 0.844).
Wiecej warstw = wiekszy receptive field = bardziej globalna agregacja = over-smoothing.
W grafach temporalnych zbyt globalna perspektywa zaciera lokalne wzorce kliniczne
(co dzialosi sie w oknie ±2h jest wazniejsze niz calosciowa srednia pobytu).

CO TO ATTENTION POOLING I DLACZEGO LEPSZY OD MEAN?
Mean pooling bierze prosta srednich cechy wszystkich nodow — traktuje pierwsza
godzine pobytu tak samo jak godzine przed kryzysu.
Attention pooling (AttentionalAggregation) uczy sie wspolczynnikow bramkujacych:
alpha_v = softmax(MLP_gate(h_v)) → g = Σ alpha_v * h_v
Model moze nauczyc sie skupiac na klinicznie waznych momentach (np. deterioracja
stanu). Wynik: +0.015 AUROC vs mean pooling (0.844 vs 0.829).

CO TO FOCAL LOSS I KIEDY LEPSZY OD WAZOWANEGO BCE?
BCE z pos_weight skaluje loss dla klasy pozytywnej stale (wspolczynnik 6.9).
Focal Loss: FL(p_t) = -alpha_t * (1-p_t)^gamma * log(p_t)
Czynnik (1-p_t)^gamma dynamicznie redukuje loss dla latwo sklasyfikowanych
przykladow (p_t duze) i skupia gradient na trudnych (p_t male).
Gamma=2.0 jest standardowym wyborem z oryginalnego paperu RetinaNet (Lin 2017).
Efekt: +0.002 AUROC, ale znaczaca poprawa sens@95spec (0.388 vs 0.326 baseline).
Klinicznie wazne: wyzszy recall przy danej specyficznosci.
""",
    # 8 — Bug / Leakage
    """\
JAK WYKRYTO LEAKAGE?
AUROC 0.931 jest podejrzanie wysoki — SOTA na tym samym zadaniu i zbiorze danych
to 0.88, a nasz poprzedni najlepszy wynik to 0.850. Kiedy wynik modelu bije
SOTA o 4%, nalezy sprawdzic preprocessing, nie swietowac.
Analizujac kod gnn_dataset.py znalezlismy: sig_rows[-max_signals:]
zamiast sig_rows[:max_signals].

DLACZEGO TO LEAKAGE?
Pacjent ktory umiera ma terminalne vitals (ekstremalne HR, BP, SpO2) bezposrednio
przed smiercia. Uzywajac OSTATNICH N sygnalow z calego pobytu model "widzial"
te sygnaly terminalne — de facto wiedzial ze pacjent zaraz umrze, zanim
zostal "zapytany" o predykcje. To nie jest predykcja, to opis.

DLACZEGO PIERWSZE N SYGNALOW?
Klinicznie sensowna predykcja: czy pacjent przyjety na OIOM umrze, na podstawie
danych z WCZESNEJ FAZY pobytu (pierwsze godziny / pierwsza doba)?
To odpowiada prawdziwemu scenariuszowi klinicznemu: lekarz chce wiedziec
o ryzyku wcze snie, nie wtedy gdy pacjent juz umiera.

CZY INNE DANE TLEZ MOGLY MIEC LEAKAGE?
Notatki kliniczne: uzywamy notatek z calego pobytu, ale sa one PRZYCZYNOWO
uzasadnione — notatka z godziny 6 pobytu NIE opisuje tego co wydarzy sie w
godzinie 48. Sygnaly all-stay po naprawieniu: pierwsze 50 = brak leakage.
ICD Charlson: kody ICD sa przypisywane przy wypisie, ale dotycza PRZEWLEKLYCH
schorzen istniejacych przed hospitalizacja — Charlson index jest wlasnie
do tego zaprojektowany.
""",
    # 9 — AUROC chart
    """\
JAK CZYTAC TEN WYKRES?
Przerywana fioletowa linia = baseline (AUROC 0.829, mean pooling bez demo i focal).
Przerywana czerwona linia = SOTA 0.880 z literatury.
Kolory pasow: czerwony = ablacje (gorsze od baseline), indigo = baseline,
zielony/cyan = warianty lepsze od baseline.

CZEMU E2E (0.790) JEST GORSZY OD BASELINE (0.829)?
Catastrophic forgetting: BERT (110M param) jest trenowany z LR ~5e-7, GNN z ~1e-3.
Roznica 200x w LR powinna byc wystarczajaca. Ale gradient z malego klasyfikatora
(65K param) przez duzy BERT (110M param) zaburza reprezentacje wyuczone w Phase 1
nawet przy tak niskim LR. Zamrozenie 8 z 12 warstw czesciowo pomaga, ale nie
eliminuje problemu. Na 8GB VRAM bez gradient checkpointing pelne E2E jest
rowniez problematyczne pamiecia.

CZEMU SIGNAL_ONLY (0.784) TAK WIELE TRACI?
Notatki kliniczne zawieraja bogate, nieskwantyfikowane informacje kliniczne —
oceny stanu ogolnego, plany leczenia, zmiany w zachowaniu pacjenta.
Sygnaly fizjologiczne sa dokladne liczbowo, ale ograniczone do mierzalnych parametrow.
Polaczenie obu: notatki daja kontekst interpretacyjny dla sygnalow.
""",
    # 10 — AUPRC chart
    """\
DLACZEGO AUPRC WAZNIEJSZA OD AUROC PRZY IMBALANSIE?
AUROC mierzy: P(score(positive) > score(negative)) dla losowej pary.
Przy 1:7 imbalansie: model moze miec AUROC 0.85 z wieloma False Negatives bo
jest bardzo wiele True Negatives ktore zawyza metryki ROC.
AUPRC skupia sie na wydajnosci dla klasy pozytywnej: precision@recall_k.
Random baseline AUPRC = prevalencja klasy = 0.127.
Random baseline AUROC = 0.5 zawsze — tak samo dla zbilansowanych i niezbilansowanych.

CZY NASZ AUROC 0.850 JEST WIARYGODNY?
Tak — mamy row niez obu metryk z podobna hierarchia eksperymentow.
Best model: AUROC 0.850, AUPRC 0.465. Relacja 0.850/0.465 vs random 0.5/0.127
jest spjna: model 3.6x powyzej losowego w precision-recall.
Gdyby AUROC byl zawyzona przez TN (False Negative problem), widzielibysmy
niska AUPRC — a mamy 0.465 czyli wyraznie powyzej 0.127.
""",
    # 11 — Learning Curve
    """\
CO POKAZUJE KRZYWA UCZENIA?
Niebieski: val AUROC na zbiorze walidacyjnym po kazdej epoce.
Szary (cienki): train loss (Focal Loss na zbiorze treningowym).
Najlepszy punkt: zaznaczony granatowym krolkiem.

CZY MODEL SIE PRZETRENOWUJE?
Val AUROC stabilizuje sie okolo 0.845 od ~40 epoki i dalej nieznacznie rosnie.
Train loss maleje monotonicznie — model dalej sie uczy, ale val AUROC nie
spada — brak overfittingu. Dropout 0.3, early stopping patience=15 zapewniaja
regularyzacje.

DLACZEGO WYNIKI 10 I 14 SA ROZNE MIMO IDENTYCZNYCH HPARAMS?
Version_10: AUROC 0.850, version_14: AUROC 0.848. Roznica 0.002 AUROC.
To jest wariancja losowa — seed deterministyczny ale PyG/CUDA operacje
mog miec niezbyt reprodukowalny wynik. Roznica mniejsza niz 0.003 jest
nieistotna statystycznie. Oznacza to ze nasz wynik jest stabilny.
""",
    # 12 — Contributions
    """\
DLACZEGO ICD NODE NIE POMAGA?
56.4% pobytow ma zerowy wektor Charlson (brak zapisanych chorob przewleklych).
Dla tych pobytow ICD node to wektor samych zer — nieinformatywny.
Pozostale 43.6% z niezerowym Charlsonem: 19-wymiarowy rzadki wektor binarny
jest za malo informatywny dla 65K-parametrowego modelu przy imbalansie 1:7.
Jesli mielibysmy wiecej parametrow lub gestszy wektor kodowania ICD, byc moze
by pomoglo. Alternatywnie: uzyc ICD jako graph-level feature (jak demografike)
zamiast osobnego nodu — to jest propozycja w future work.

DLACZEGO ALL-STAY SIGNALS (0.820) GORSZE OD BASELINE (0.829)?
Baseline uzywa sygnalow w oknie ±2h od kazdej notatki — to okno jest juz
wybrane przez Phase 1 jako klinicznie relevantne. Wziecie pierwszych 50
sygnalow z calego pobytu (medianicznie 315 sygnalow, max 30K) przy cap=50
daje fragmentaryczne pokrycie calego pobytu. Okno ±2h jest bogato pokryte
przez Phase 1; random sample całego pobytu jest zbyt rzadki.
Potencjalne rozwiazanie: wiecej sygnałow (cap=200+), ale wymaga wiecej RAM.

DLACZEGO FOCAL LOSS BARDZIEJ POMAGA NA SENS@95SPEC NIZ AUROC?
Focal Loss skupia gradient na trudnych przykladach — czyli pacjentach
trudno klasyfikowalnych (ci co sa blisko granicy decyzyjnej). Przy wysokiej
specyficznosci (95%) chcemy wylapac tak duzo True Positives jak mozliwe.
Focal Loss przesuwa gradient w kierunku wlasciwej klasyfikacji trudnych
przypadkow pozytywnych — stad poprawa recallu przy danym progu specyficznosci.
Nie zmienia dramatycznie AUROC (globalny ranking), ale poprawia kliniczne
punkt decyzyjny.
""",
    # 13 — Best Model
    """\
CZEMU AKURAT ATTENTION A NIE DUAL POOLING?
Dual pooling (osobne pule dla nodow sygnalow i notatek, konkatenowane) daje
AUROC 0.846 — prawie tak dobry jak attention (0.844 bez focal, 0.850 z focal).
Dual pooling ma lepsza kalibracje (Brier score 0.170 vs 0.195 dla attention).
Attention pooling + focal loss daje najlepsze AUROC i sens@95spec.
W zastosowaniach klinicznych gdzie zalezy nam na kalibracji (precyzja prawdopod.)
dual pooling moze byc lepszym wyborem. Dla maksymalizacji AUROC — attention+focal.

CO OZNACZA SENS@95SPEC = 0.388?
Przy ustawieniu progu klasyfikacji tak ze 95% ocalonych jest poprawnie
identyfikowanych jako "niskie ryzyko" (specificity=0.95) — model prawidlowo
identyfikuje 38.8% pacjentow wysokiego ryzyka (sensitivity=0.388).
W praktyce: ze 100 pacjentow smiertelnych na OIOM model oznaczy alarmem ~39 z nich
przy tylko 5% falfyw alarmow wsrod ocalonych. To jest klinicznie uzyteczne —
szczegolnie w srodowisku z ograniczonymi zasobami.

DLACZEGO REDUCELRONPLATEAU A NIE COSINE DECAY?
ReduceLROnPlateau monitoruje val_auroc i zmniejsza LR o polowe gdy brak
poprawy przez 5 epok (patience=5, factor=0.5). Jest adaptacyjny — nie wymaga
ustalonego harmonogramu. Cosine decay bylby odpowiedni gdybysmy znali dokladna
liczbe epok z gory; tutaj early stopping z patience=15 przerywa trening
w roznych momentach dla roznych eksperymentow.
""",
    # 14 — vs SOTA
    """\
CO DOKLADNIE ROBI SOTA KTORE MY NIE ROBIMY?
Modele SOTA na predykcji smiertelnosci (np. ClinicalBERT + temporal attention,
MIMIC-Extract + transformer) uzywaja:
1. End-to-end fine-tuning BERT z gradient checkpointing (eliminuje OOM)
2. Multi-GPU distributed training (mozna trenowac batch_size=64+)
3. Temporal attention zamiast GINEConv (lepsze modelowanie nieregularnych sygnalow)
4. Czesto wiecej danych pomocniczych (leki, procedury) ktorych my nie uwzgledniamy

DLACZEGO NIE MOZNA BYLO TEGO ZROBIC NA NASZYM SPRZECIE?
RTX 4060 Laptop GPU: 8GB VRAM.
Pelny E2E BERT (110M param) + batch=32 = ~12GB VRAM (OOM).
Gradient checkpointing zmniejszyloby zu zycie 4x kosztem 30% spowolnienia —
nie zaimplementowane w obecnej wersji. Na dedykowanym serwerze (A100 80GB) E2E
byloby mozliwe w ~5-8 min/epoke zamiast crash.

JAK MIERZYC WKLAD NASZEJ PRACY SKORO NIE BILIMY SOTA?
1. Systematyczne ablacje: pokazujemy co konkretnie wnosi kazda skladowa
2. Wykrycie i naprawa temporal leakage: metodologiczny wklad
3. Dokumentacja negatywnych wynikow: wartosc dla powtarzalnosci nauki
4. Efektywnosc parametrow: 65K param osiaga 96.6% wydajnosci SOTA (110M param)
5. Otwarta reprodukowalna implementacja na MIMIC-IV 3.1
""",
    # 15 — Takeaways
    """\
PODSUMOWANIE KLUCZOWYCH DECYZJI PROJEKTOWYCH:

1. GINECONV: jedyny standardowy GNN oprocz PNA ktory natywnie obsluguje cechy
   krawedzi. Delta_t jest kluczowe dla temporalnego kontekstu klinicznego.

2. PHASE 1 ZAMIAST E2E: catastrophic forgetting potwierdza ze rzad wielkosci
   roznica LR (1e-3 vs 5e-7) nie wystarczy. Zamrozenie reprezentacji zachowuje
   wiedze zdobyta na 102K notatkach.

3. INFONCE: efektywny batch-level contrastive loss. O(B^2) negatywnych par
   vs O(B) w triplet loss przy tym samym koszcie obliczeniowym.

4. TEMPORAL LEAKAGE: kluczowa lekcja metodologiczna. W danych medycznych
   sekwencyjnych zawsze sprawdzaj kierunek czasu w preprocessingu.

5. FOCAL LOSS: poprawia klinicznie istotna metryki (sens@95spec) nawet jesli
   AUROC rosnie marginalnie. Wazne jesli model ma byc uzyteczny w praktyce.

CZY PROJEKT JEST POWTARZALNY?
Tak. Kod jest publiczny, dane dostepne przez PhysioNet po rejestracji.
Seed deterministyczny (42). Szczegolowe logi w data/snapshots/gnn/.
Jedyna potencjalna niereprodukowalnosc: niedeterministyczne operacje CUDA.
Odchylenie standardowe miedzy uruchomieniami ~0.002 AUROC.
""",
]


def _add_notes(slide, text: str) -> None:
    """Dodaje notatki do slajdu (widoczne w PowerPoint Presenter View)."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.clear()
    for i, para_text in enumerate(text.strip().splitlines()):
        if i == 0:
            tf.paragraphs[0].text = para_text
        else:
            p = tf.add_paragraph()
            p.text = para_text


if __name__ == "__main__":
    main()
